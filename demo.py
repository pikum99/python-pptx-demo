"""
python-pptxでPowerPoint生成のデモ
- テーブル
- グラフ（棒グラフ、折れ線グラフ、円グラフ、帯グラフ）
- テキストボックス
- 箇条書き
- スタイリング
"""

from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# スタイル定数
TITLE_FONT_SIZE = Pt(32)
CHART_Y_POSITION = Inches(1.2)
CHART_HEIGHT = Inches(4.2)
BRAND_COLOR = RGBColor(0, 51, 102)

def add_title_slide(prs, title, subtitle):
    """タイトルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = BRAND_COLOR

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

def add_table_slide(prs, data):
    """テーブルスライドを追加"""
    if not data:
        raise ValueError("データが空です")

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 完全な空白レイアウト

    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = "データテーブル"
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = BRAND_COLOR

    # テーブル作成
    rows = len(data) + 1
    cols = len(data[0])
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(0.5 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # ヘッダー行のスタイリング
    for i, key in enumerate(data[0].keys()):
        cell = table.cell(0, i)
        cell.text = key
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_COLOR
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # データ行
    for row_idx, item in enumerate(data, start=1):
        for col_idx, value in enumerate(item.values()):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)

            # 交互に背景色を変える
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

def add_bar_chart_slide(prs, data):
    """棒グラフスライドを追加"""
    if not data:
        raise ValueError("データが空です")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = "売上比較（棒グラフ）"
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True

    # チャートデータ準備
    chart_data = CategoryChartData()
    chart_data.categories = [item['製品'] for item in data]
    chart_data.add_series('売上', [item['売上'] for item in data])

    # チャート追加
    x, y, cx, cy = Inches(1.5), CHART_Y_POSITION, Inches(7), CHART_HEIGHT
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # チャートスタイリング
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 軸ラベルの設定
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = True

def add_line_chart_slide(prs, time_data):
    """折れ線グラフスライドを追加"""
    if not time_data:
        raise ValueError("データが空です")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = "月次推移（折れ線グラフ）"
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True

    # チャートデータ準備
    chart_data = CategoryChartData()
    chart_data.categories = [item['月'] for item in time_data]
    chart_data.add_series('売上', [item['売上'] for item in time_data])
    chart_data.add_series('利益', [item['利益'] for item in time_data])

    # チャート追加
    x, y, cx, cy = Inches(1.5), CHART_Y_POSITION, Inches(7), CHART_HEIGHT
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    # チャートスタイリング
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 軸ラベルの設定
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = True

def add_stacked_bar_chart_slide(prs):
    """帯グラフ（積み上げ横棒グラフ）スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = "評価比較（帯グラフ）"
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True

    # チャートデータ準備
    chart_data = CategoryChartData()
    chart_data.categories = ['上位10', '上位100']
    chart_data.add_series('Good', (30, 50))
    chart_data.add_series('Bad', (70, 50))

    # チャート追加（横向き100%積み上げ棒グラフ）
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(3.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data
    ).chart

    # チャートスタイリング
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

def add_stacked_column_chart_slide(prs):
    """縦向き積み上げ100%棒グラフスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = "評価比較（縦向き積み上げグラフ）"
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True

    # チャートデータ準備
    chart_data = CategoryChartData()
    chart_data.categories = ['上位10', '上位100']
    chart_data.add_series('Good', (30, 50))
    chart_data.add_series('Bad', (70, 50))

    # チャート追加（縦向き100%積み上げ棒グラフ）
    x, y, cx, cy = Inches(2.5), CHART_Y_POSITION, Inches(5), CHART_HEIGHT
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED_100, x, y, cx, cy, chart_data
    ).chart

    # チャートスタイリング
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

def add_pie_chart_slide(prs, data):
    """円グラフスライドを追加"""
    if not data:
        raise ValueError("データが空です")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = "市場シェア（円グラフ）"
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True

    # チャートデータ準備
    chart_data = CategoryChartData()
    chart_data.categories = [item['製品'] for item in data]
    chart_data.add_series('シェア', [item['売上'] for item in data])

    # チャート追加
    x, y, cx, cy = Inches(2), CHART_Y_POSITION, Inches(6), CHART_HEIGHT
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    # チャートスタイリング
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

def add_bullet_slide(prs, title, bullet_points):
    """箇条書きスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # タイトル＋コンテンツ

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_shape.text_frame.paragraphs[0].font.bold = True

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()  # 既存のテキストをクリア

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # サブポイントがある場合
        if isinstance(point, dict):
            p.text = point['main']
            p.level = 0
            p.font.size = Pt(18)
            for sub_point in point.get('sub', []):
                p = text_frame.add_paragraph()
                p.text = sub_point
                p.level = 1
                p.font.size = Pt(16)
        else:
            p.text = point
            p.level = 0
            p.font.size = Pt(18)

def add_two_column_slide(prs, title, left_content, right_content):
    """2カラムレイアウトのスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True

    # 左カラム
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "メリット"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_before = Pt(6)

    # 右カラム
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "デメリット"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)

    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_before = Pt(6)

def add_summary_slide(prs):
    """まとめスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景色を設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLOR

    # 中央にテキスト
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    text_frame = textbox.text_frame
    text_frame.text = "ご清聴ありがとうございました"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def main():
    """メイン処理"""
    try:
        # プレゼンテーション作成
        prs = Presentation()
        prs.slide_width = Inches(10)   # 16:9
        prs.slide_height = Inches(5.625)

        # サンプルデータ
        product_data = [
            {"製品": "製品A", "売上": 1000, "利益": 200, "市場シェア": "25%"},
            {"製品": "製品B", "売上": 1500, "利益": 350, "市場シェア": "35%"},
            {"製品": "製品C", "売上": 800, "利益": 150, "市場シェア": "20%"},
            {"製品": "製品D", "売上": 600, "利益": 100, "市場シェア": "15%"},
            {"製品": "その他", "売上": 200, "利益": 30, "市場シェア": "5%"},
        ]

        time_series_data = [
            {"月": "1月", "売上": 800, "利益": 150},
            {"月": "2月", "売上": 900, "利益": 180},
            {"月": "3月", "売上": 1100, "利益": 220},
            {"月": "4月", "売上": 1000, "利益": 200},
            {"月": "5月", "売上": 1200, "利益": 250},
            {"月": "6月", "売上": 1400, "利益": 300},
        ]

        # スライド1: タイトル
        add_title_slide(
            prs,
            "2024年度 ビジネスレポート",
            f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"
        )

        # スライド2: 箇条書き（目次）
        add_bullet_slide(
            prs,
            "本日のアジェンダ",
            [
                "データテーブルの概要",
                "売上比較分析",
                "月次推移の確認",
                "市場シェア分析",
                "まとめと今後の展望"
            ]
        )

        # スライド3: テーブル
        add_table_slide(prs, product_data)

        # スライド4: 棒グラフ
        add_bar_chart_slide(prs, product_data)

        # スライド5: 折れ線グラフ
        add_line_chart_slide(prs, time_series_data)

        # スライド6: 円グラフ
        add_pie_chart_slide(prs, product_data)

        # スライド6.5: 帯グラフ（横向き）
        add_stacked_bar_chart_slide(prs)

        # スライド6.6: 縦向き積み上げグラフ
        add_stacked_column_chart_slide(prs)

        # スライド7: 2カラムレイアウト
        add_two_column_slide(
            prs,
            "製品Bの評価",
            left_content=[
                "売上トップ",
                "高い利益率",
                "市場シェア拡大中",
                "顧客満足度が高い"
            ],
            right_content=[
                "競合の増加",
                "原価上昇リスク",
                "季節変動が大きい"
            ]
        )

        # スライド8: 階層的な箇条書き
        add_bullet_slide(
            prs,
            "今後のアクションプラン",
            [
                {
                    "main": "製品Bの強化",
                    "sub": ["マーケティング予算増額", "新規チャネル開拓"]
                },
                {
                    "main": "製品Cの改善",
                    "sub": ["コスト削減施策", "品質向上プログラム"]
                },
                "競合分析の強化",
                "四半期レビューの実施"
            ]
        )

        # スライド9: まとめ
        add_summary_slide(prs)

        # ファイル保存
        output_path = "demo_presentation.pptx"
        prs.save(output_path)
        print(f"プレゼンテーション作成完了: {output_path}")

        return output_path

    except Exception as e:
        print(f"エラーが発生しました: {e}")
        raise

if __name__ == "__main__":
    main()